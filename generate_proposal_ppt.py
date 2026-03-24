#!/usr/bin/env python3
"""마드모아젤헤어 플레이스 순위 회복 제안서 PPT 생성 스크립트"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ── 컬러 시스템 ──
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0xF8, 0xF9, 0xFA)
C_DARK_NAVY = RGBColor(0x1B, 0x2A, 0x4A)
C_BLUE = RGBColor(0x25, 0x63, 0xEB)
C_GOLD = RGBColor(0xD4, 0xA8, 0x53)
C_DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
C_TABLE_ALT1 = RGBColor(0xF0, 0xF4, 0xFF)
C_TABLE_ALT2 = RGBColor(0xFF, 0xFF, 0xFF)
C_PLACEHOLDER_BG = RGBColor(0xE5, 0xE7, 0xEB)
C_PLACEHOLDER_BORDER = RGBColor(0x9C, 0xA3, 0xAF)
C_GREEN = RGBColor(0x16, 0xA3, 0x4A)
C_GOLD_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
C_RED = RGBColor(0xDC, 0x26, 0x26)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_NAME = "맑은 고딕"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ── 유틸리티 함수 ──

def add_bg(slide, color):
    """슬라이드 배경색 설정"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None,
              line_width=None, shape_type=MSO_SHAPE.RECTANGLE):
    """도형 추가"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.fill.solid()
        shape.line.width = line_width or Pt(1)
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=C_DARK_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_NAME, line_spacing=1.2):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, default_size=14,
                       default_color=C_DARK_TEXT, default_bold=False,
                       alignment=PP_ALIGN.LEFT, line_spacing=1.5):
    """여러 줄 텍스트 박스. lines = [(text, size, bold, color), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, sz, bld, clr = line, default_size, default_bold, default_color
        else:
            txt = line[0]
            sz = line[1] if len(line) > 1 else default_size
            bld = line[2] if len(line) > 2 else default_bold
            clr = line[3] if len(line) > 3 else default_color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.bold = bld
        p.font.color.rgb = clr
        p.font.name = FONT_NAME
        p.alignment = alignment
        p.space_after = Pt(4)
        p.line_spacing = Pt(sz * line_spacing)
    return txBox


def add_top_line(slide, color=C_BLUE, height=Pt(4)):
    """상단 블루 라인"""
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, height, fill_color=color)


def add_footer(slide, page_num, text_color=C_DARK_TEXT):
    """하단 푸터"""
    add_text_box(slide, Inches(10.5), Inches(6.95), Inches(2.5), Inches(0.4),
                 f"온다마케팅  |  {page_num}", font_size=9, color=text_color,
                 alignment=PP_ALIGN.RIGHT)


def add_accent_box(slide, left, top, width, height, text, border_color=C_BLUE,
                   bg_color=None, text_color=C_DARK_TEXT, font_size=13, bold=False):
    """강조 박스 (테두리 + 텍스트)"""
    box = add_shape(slide, left, top, width, height,
                    fill_color=bg_color or C_WHITE,
                    line_color=border_color, line_width=Pt(2))
    txBox = add_text_box(slide, left + Inches(0.3), top + Inches(0.15),
                         width - Inches(0.6), height - Inches(0.3),
                         text, font_size=font_size, bold=bold, color=text_color)
    return box


def add_image_placeholder(slide, left, top, width, height, label="이미지"):
    """이미지 플레이스홀더"""
    box = add_shape(slide, left, top, width, height,
                    fill_color=C_PLACEHOLDER_BG,
                    line_color=C_PLACEHOLDER_BORDER, line_width=Pt(2))
    # 대시 스타일은 python-pptx에서 제한적이므로 실선 사용
    add_text_box(slide, left, top + height / 2 - Inches(0.2), width, Inches(0.4),
                 f"[이미지 자리: {label}]", font_size=11,
                 color=C_PLACEHOLDER_BORDER, alignment=PP_ALIGN.CENTER)
    return box


def set_cell_style(cell, text, font_size=11, bold=False, color=C_DARK_TEXT,
                   fill_color=None, alignment=PP_ALIGN.CENTER):
    """테이블 셀 스타일"""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color


def add_icon_circle(slide, left, top, size, emoji, bg_color=C_BLUE):
    """원형 아이콘 + 이모지"""
    circle = add_shape(slide, left, top, size, size,
                       fill_color=bg_color, shape_type=MSO_SHAPE.OVAL)
    add_text_box(slide, left, top + size * 0.15, size, Inches(0.5),
                 emoji, font_size=24, alignment=PP_ALIGN.CENTER, color=C_WHITE)
    return circle


def add_card(slide, left, top, width, height, title, subtitle, detail,
             bg_color=C_WHITE, border_color=C_BLUE):
    """카드 도형"""
    card = add_shape(slide, left, top, width, height,
                     fill_color=bg_color, line_color=border_color, line_width=Pt(1.5))
    y = top + Inches(0.25)
    add_text_box(slide, left + Inches(0.25), y, width - Inches(0.5), Inches(0.35),
                 title, font_size=14, bold=True, color=C_DARK_TEXT,
                 alignment=PP_ALIGN.CENTER)
    y += Inches(0.45)
    add_text_box(slide, left + Inches(0.25), y, width - Inches(0.5), Inches(0.3),
                 subtitle, font_size=12, color=C_BLUE, alignment=PP_ALIGN.CENTER)
    y += Inches(0.35)
    add_text_box(slide, left + Inches(0.25), y, width - Inches(0.5), Inches(0.25),
                 detail, font_size=11, color=C_GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER)
    return card


def new_slide():
    """빈 슬라이드 추가"""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def make_table(slide, left, top, width, rows, cols, row_height=Inches(0.45)):
    """테이블 생성"""
    tbl_h = row_height * rows
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, tbl_h)
    table = table_shape.table
    for i in range(rows):
        table.rows[i].height = row_height
    return table


# ── 슬라이드 생성 ──

def slide_01_cover():
    """표지"""
    slide = new_slide()
    add_bg(slide, C_DARK_NAVY)
    # 상단 골드 라인
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Pt(4), fill_color=C_GOLD)
    # 상단 소제목
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
                 "네이버 종합광고 실행사", font_size=14, color=C_GOLD,
                 alignment=PP_ALIGN.CENTER)
    # 중앙 대형 타이틀
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
                 "마드모아젤헤어", font_size=44, bold=True, color=C_WHITE,
                 alignment=PP_ALIGN.CENTER)
    # 서브타이틀
    add_text_box(slide, Inches(1), Inches(3.7), Inches(11), Inches(0.7),
                 "플레이스 순위 회복 및 방어 전략 제안서", font_size=24,
                 bold=True, color=C_GOLD, alignment=PP_ALIGN.CENTER)
    # 설명
    add_text_box(slide, Inches(1), Inches(4.7), Inches(11), Inches(0.5),
                 "세종시 4개 지점 통합 진단  |  안전보장형 운영 제안", font_size=14,
                 color=RGBColor(0xA0, 0xAE, 0xC9), alignment=PP_ALIGN.CENTER)
    # 하단 구분선
    add_shape(slide, Inches(5), Inches(5.5), Inches(3), Pt(1), fill_color=C_GOLD)
    # 회사명 + 날짜
    add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
                 "온다마케팅  |  2026.03", font_size=12,
                 color=RGBColor(0x80, 0x90, 0xB0), alignment=PP_ALIGN.CENTER)


def slide_02_greeting():
    """인사말"""
    slide = new_slide()
    add_bg(slide, C_WHITE)
    add_top_line(slide)
    add_footer(slide, 2)
    # 제목
    add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                 "인사말", font_size=32, bold=True, color=C_DARK_NAVY)
    # 구분선
    add_shape(slide, Inches(1.2), Inches(1.3), Inches(2), Pt(3), fill_color=C_BLUE)
    # 본문
    body = (
        "안녕하세요, 온다마케팅입니다.\n\n"
        "저희는 네이버 플레이스, 인스타그램, 블로그 등 온라인 채널 통합 운영을 전문으로 하는\n"
        "종합광고 실행사입니다.\n\n"
        "단순한 대행이 아닌, 데이터 기반의 전략 설계와 실행력으로\n"
        "광고주의 실질적인 매출 성장을 함께 만들어가고 있습니다.\n\n"
        "마드모아젤헤어 대표님께 감사드리며,\n"
        "본 제안서가 귀사의 성장에 도움이 되기를 바랍니다."
    )
    add_text_box(slide, Inches(1.2), Inches(1.7), Inches(10), Inches(3.5),
                 body, font_size=14, color=C_DARK_TEXT, line_spacing=1.6)
    # 하단 강조 박스
    add_accent_box(slide, Inches(1.2), Inches(5.3), Inches(11), Inches(1),
                   "본 제안서는 귀사의 현재 플레이스 현황에 대한 진단과 함께,\n"
                   "2026년 변화된 알고리즘 환경에 맞는 대응 전략을 담고 있습니다.",
                   border_color=C_BLUE, bg_color=RGBColor(0xEF, 0xF6, 0xFF),
                   font_size=13, bold=False)


def slide_03_rule_change():
    """2026년 규칙 변화"""
    slide = new_slide()
    add_bg(slide, C_LIGHT_GRAY)
    add_top_line(slide)
    add_footer(slide, 3)
    # 제목
    add_text_box(slide, Inches(1.2), Inches(0.6), Inches(11), Inches(0.6),
                 "2026년, 플레이스 마케팅의 규칙이 바뀌었습니다",
                 font_size=28, bold=True, color=C_DARK_NAVY)
    add_shape(slide, Inches(1.2), Inches(1.3), Inches(2), Pt(3), fill_color=C_BLUE)

    # 좌측: 과거
    left_x = Inches(1.2)
    box_w = Inches(5.2)
    add_shape(slide, left_x, Inches(1.7), box_w, Inches(3.3),
              fill_color=C_WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
    add_text_box(slide, left_x, Inches(1.75), box_w, Inches(0.45),
                 "  과거 (2025년 이전)", font_size=16, bold=True, color=C_WHITE)
    add_shape(slide, left_x, Inches(1.7), box_w, Inches(0.45),
              fill_color=RGBColor(0x6B, 0x72, 0x80))
    add_text_box(slide, left_x + Inches(0.05), Inches(1.75), box_w - Inches(0.1), Inches(0.45),
                 "과거 (2025년 이전)", font_size=16, bold=True, color=C_WHITE,
                 alignment=PP_ALIGN.CENTER)
    past_items = [
        "트래픽 양 = 순위 (단순 공식)",
        "3~6개월 단위 로직 변동",
        "배포형 블로그로 손쉽게 점수 확보",
        "공격적 운영 = 빠른 성과"
    ]
    for i, item in enumerate(past_items):
        add_text_box(slide, left_x + Inches(0.4), Inches(2.3 + i * 0.55),
                     box_w - Inches(0.8), Inches(0.4),
                     f"▸  {item}", font_size=13, color=C_DARK_TEXT)

    # 화살표
    add_text_box(slide, Inches(6.3), Inches(2.8), Inches(0.8), Inches(0.6),
                 "→", font_size=36, bold=True, color=C_BLUE,
                 alignment=PP_ALIGN.CENTER)

    # 우측: 현재
    right_x = Inches(6.9)
    add_shape(slide, right_x, Inches(1.7), box_w, Inches(3.3),
              fill_color=C_WHITE, line_color=C_BLUE, line_width=Pt(2))
    add_shape(slide, right_x, Inches(1.7), box_w, Inches(0.45),
              fill_color=C_BLUE)
    add_text_box(slide, right_x + Inches(0.05), Inches(1.75), box_w - Inches(0.1), Inches(0.45),
                 "현재 (2026년 이후)", font_size=16, bold=True, color=C_WHITE,
                 alignment=PP_ALIGN.CENTER)
    current_items = [
        "패턴 감지 → 즉시 제재",
        "1~2주 단위 가중치 급변동",
        "블로그 대량 삭제 + 순위 급락",
        "공격적 운영 = 공격적 하락"
    ]
    for i, item in enumerate(current_items):
        add_text_box(slide, right_x + Inches(0.4), Inches(2.3 + i * 0.55),
                     box_w - Inches(0.8), Inches(0.4),
                     f"▸  {item}", font_size=13, bold=True, color=C_RED)

    # 하단 강조 박스
    add_accent_box(slide, Inches(1.2), Inches(5.4), Inches(11), Inches(0.8),
                   "핵심:  순위를 '올리는 것'보다, '지키는 것'이 더 어렵고 더 중요한 시대",
                   border_color=C_GOLD, bg_color=RGBColor(0xFF, 0xFA, 0xEB),
                   text_color=C_DARK_NAVY, font_size=15, bold=True)


def slide_04_defense_system():
    """온다마케팅 대응 체계"""
    slide = new_slide()
    add_bg(slide, C_WHITE)
    add_top_line(slide)
    add_footer(slide, 4)
    add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                 "데이터 기반 순위 방어 체계", font_size=30, bold=True, color=C_DARK_NAVY)
    add_shape(slide, Inches(1.2), Inches(1.3), Inches(2), Pt(3), fill_color=C_BLUE)

    # 3개 원형 아이콘 가로 배치
    items = [
        ("🔬", "셀프마케팅 + AI로직연구소", "구축 완료", C_BLUE),
        ("🛡️", "구체화된 순위 방어 로직", "+ 상승 로직 적용", RGBColor(0x7C, 0x3A, 0xED)),
        ("📈", "안전하고 꾸준한 상승", "데이터 기반 추구", C_GREEN),
    ]
    for i, (emoji, title, sub, clr) in enumerate(items):
        cx = Inches(1.5 + i * 4)
        # 배경 카드
        add_shape(slide, cx, Inches(1.8), Inches(3.3), Inches(3.2),
                  fill_color=C_LIGHT_GRAY, line_color=RGBColor(0xE0, 0xE0, 0xE0),
                  line_width=Pt(1))
        # 아이콘 원
        add_icon_circle(slide, cx + Inches(1.15), Inches(2.2), Inches(1), emoji, bg_color=clr)
        # 텍스트
        add_text_box(slide, cx + Inches(0.15), Inches(3.5), Inches(3), Inches(0.5),
                     title, font_size=15, bold=True, color=C_DARK_NAVY,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, cx + Inches(0.15), Inches(4.0), Inches(3), Inches(0.4),
                     sub, font_size=12, color=clr, alignment=PP_ALIGN.CENTER)

    # 하단 설명
    add_text_box(slide, Inches(1.2), Inches(5.6), Inches(11), Inches(0.5),
                 "단순 트래픽 유입 대행에서 벗어나, 실시간 데이터 분석 기반 순위 방어 체계 구축",
                 font_size=14, color=RGBColor(0x64, 0x64, 0x64),
                 alignment=PP_ALIGN.CENTER)


def slide_05_diagnosis():
    """현황 진단"""
    slide = new_slide()
    add_bg(slide, C_WHITE)
    add_top_line(slide)
    add_footer(slide, 5)
    add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                 "마드모아젤헤어 4개 지점 현황 진단",
                 font_size=28, bold=True, color=C_DARK_NAVY)
    add_shape(slide, Inches(1.2), Inches(1.3), Inches(2), Pt(3), fill_color=C_BLUE)

    # 테이블
    headers = ["지점", "순위 변동", "N2 변동", "블로그 삭제", "저품질"]
    data = [
        ["마크원애비뉴점", "하락", "-0.08 이상", "해당", "없음"],
        ["나성1호점", "1위 → 14위", "-0.08 이상", "해당", "없음"],
        ["나성2호점", "14위 → 21위", "-0.10", "해당", "없음"],
        ["반곡점", "하락", "-0.08 이상", "해당", "없음"],
    ]
    tbl = make_table(slide, Inches(1.2), Inches(1.7), Inches(11), 5, 5, Inches(0.5))
    col_widths = [Inches(2.5), Inches(2.2), Inches(2.2), Inches(2.1), Inches(2)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w
    for j, h in enumerate(headers):
        set_cell_style(tbl.cell(0, j), h, font_size=12, bold=True,
                       color=C_WHITE, fill_color=C_DARK_NAVY)
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            bg = C_TABLE_ALT1 if r % 2 == 0 else C_TABLE_ALT2
            clr = C_RED if val in ["하락", "해당"] or "→" in val else C_DARK_TEXT
            if val == "없음":
                clr = C_GREEN
            set_cell_style(tbl.cell(r + 1, c), val, font_size=11,
                           color=clr, fill_color=bg)

    # 공통사항
    add_multiline_text(slide, Inches(1.2), Inches(4.3), Inches(11), Inches(1),
                       [
                           ("공통사항", 13, True, C_DARK_NAVY),
                           ("• N1(SEO) 변동 없음 — 플레이스 자체 문제 아님", 12, False, C_DARK_TEXT),
                           ("• 배포형 블로그 300~400건 일괄 삭제로 N2 점수 급락", 12, False, C_DARK_TEXT),
                       ])

    # 강조 박스
    add_accent_box(slide, Inches(1.2), Inches(5.7), Inches(11), Inches(0.8),
                   "✅  종합 판단:  4개 지점 모두 저품질 상태 아님  →  회복 가능",
                   border_color=C_GREEN, bg_color=RGBColor(0xEC, 0xFD, 0xF5),
                   text_color=C_DARK_NAVY, font_size=16, bold=True)


def slide_06_qa_n2():
    """Q&A - N2 동시 하락 원인"""
    slide = new_slide()
    add_bg(slide, C_LIGHT_GRAY)
    add_top_line(slide)
    add_footer(slide, 6)

    # Q
    add_shape(slide, Inches(1.2), Inches(0.7), Inches(0.7), Inches(0.7),
              fill_color=C_BLUE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(1.2), Inches(0.78), Inches(0.7), Inches(0.5),
                 "Q", font_size=28, bold=True, color=C_WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.1), Inches(0.78), Inches(10), Inches(0.6),
                 "4개 지점 N2가 동시에 하락한 원인은 무엇입니까?",
                 font_size=20, bold=True, color=C_DARK_NAVY)

    # A
    add_shape(slide, Inches(1.2), Inches(1.8), Inches(0.7), Inches(0.7),
              fill_color=C_GOLD, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(1.2), Inches(1.88), Inches(0.7), Inches(0.5),
                 "A", font_size=28, bold=True, color=C_WHITE,
                 alignment=PP_ALIGN.CENTER)

    answer = (
        "2026년 3월 초, 네이버가 배포형 블로그를 대량 삭제하면서\n"
        "4개 지점 모두에서 N2(외부 유입 점수)가 동시에 하락했습니다.\n\n"
        "이는 특정 지점만의 문제가 아닌, 네이버 알고리즘 정책 변화에 따른\n"
        "업계 전반적인 현상입니다."
    )
    add_text_box(slide, Inches(2.1), Inches(1.85), Inches(10), Inches(2),
                 answer, font_size=14, color=C_DARK_TEXT, line_spacing=1.6)

    # 핵심 3줄 강조
    highlights = [
        ("✅  N1(SEO) 변동 없음", C_GREEN),
        ("✅  4개 지점 모두 저품질 해당 없음", C_GREEN),
        ("✅  정상적인 회복 가능", C_BLUE),
    ]
    for i, (txt, clr) in enumerate(highlights):
        y = Inches(4.2 + i * 0.6)
        add_shape(slide, Inches(2.1), y, Inches(8), Inches(0.5),
                  fill_color=C_WHITE, line_color=clr, line_width=Pt(1.5))
        add_text_box(slide, Inches(2.4), y + Inches(0.05), Inches(7.5), Inches(0.4),
                     txt, font_size=14, bold=True, color=clr)


def slide_07_qa_blog():
    """Q&A - 블로그 삭제 상쇄"""
    slide = new_slide()
    add_bg(slide, C_WHITE)
    add_top_line(slide)
    add_footer(slide, 7)

    # Q
    add_shape(slide, Inches(1.2), Inches(0.7), Inches(0.7), Inches(0.7),
              fill_color=C_BLUE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(1.2), Inches(0.78), Inches(0.7), Inches(0.5),
                 "Q", font_size=28, bold=True, color=C_WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.1), Inches(0.78), Inches(10), Inches(0.6),
                 "블로그 대량 삭제, 어떻게 상쇄합니까?",
                 font_size=20, bold=True, color=C_DARK_NAVY)

    # A
    add_shape(slide, Inches(1.2), Inches(1.8), Inches(0.7), Inches(0.7),
              fill_color=C_GOLD, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(1.2), Inches(1.88), Inches(0.7), Inches(0.5),
                 "A", font_size=28, bold=True, color=C_WHITE,
                 alignment=PP_ALIGN.CENTER)

    add_multiline_text(slide, Inches(2.1), Inches(1.85), Inches(10), Inches(3.5),
                       [
                           ("단순 재배포 = 오히려 독", 16, True, C_RED),
                           ("", 8),
                           ("CLOVA AI 도입 이후, 동일 패턴의 블로그를 재배포하면", 13, False, C_DARK_TEXT),
                           ("오히려 패턴으로 인식되어 추가 제재를 받습니다.", 13, False, C_DARK_TEXT),
                           ("", 8),
                           ("핵심 포인트", 14, True, C_DARK_NAVY),
                           ("• 2주 단위, 빠르면 1주 단위 가중치 변동 환경", 13, False, C_DARK_TEXT),
                           ("• 패턴화되지 않는 다양한 경로의 유입 설계 필요", 13, False, C_DARK_TEXT),
                           ("• AI가 감지할 수 없는 자연스러운 시그널 구축", 13, False, C_DARK_TEXT),
                       ])

    # 강조 박스
    add_accent_box(slide, Inches(1.2), Inches(5.5), Inches(11), Inches(0.8),
                   "과거 방식의 '양'으로 밀어붙이는 전략은 2026년에 더 이상 통하지 않습니다.\n"
                   "패턴을 피하면서 점수를 회복하는 정밀한 전략이 필요합니다.",
                   border_color=C_BLUE, font_size=13, bold=True)


def slide_08_qa_recovery():
    """Q&A - 회복 플랜"""
    slide = new_slide()
    add_bg(slide, C_LIGHT_GRAY)
    add_top_line(slide)
    add_footer(slide, 8)

    # Q
    add_shape(slide, Inches(1.2), Inches(0.7), Inches(0.7), Inches(0.7),
              fill_color=C_BLUE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(1.2), Inches(0.78), Inches(0.7), Inches(0.5),
                 "Q", font_size=28, bold=True, color=C_WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.1), Inches(0.78), Inches(10), Inches(0.6),
                 "급락 지점의 1단계 회복 플랜은?",
                 font_size=20, bold=True, color=C_DARK_NAVY)

    # A
    add_shape(slide, Inches(1.2), Inches(1.8), Inches(0.7), Inches(0.7),
              fill_color=C_GOLD, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(1.2), Inches(1.88), Inches(0.7), Inches(0.5),
                 "A", font_size=28, bold=True, color=C_WHITE,
                 alignment=PP_ALIGN.CENTER)

    add_multiline_text(slide, Inches(2.1), Inches(1.85), Inches(10), Inches(3),
                       [
                           ("패턴화되지 않는 작업 + AI 감지 불가 방식 설계", 15, True, C_DARK_NAVY),
                           ("", 8),
                           ("1단계 회복 전략 핵심:", 13, True, C_BLUE),
                           ("• 기존 삭제된 블로그와 다른 패턴의 콘텐츠 설계", 13, False, C_DARK_TEXT),
                           ("• 다양한 유입 경로 분산 (블로그 + 카페 + 지도저장 등)", 13, False, C_DARK_TEXT),
                           ("• 급격한 유입 증가 없이 자연스러운 회복 곡선 설계", 13, False, C_DARK_TEXT),
                           ("• 1~2주 단위 모니터링 + 실시간 전략 조정", 13, False, C_DARK_TEXT),
                       ])

    # 약속 박스
    add_accent_box(slide, Inches(1.2), Inches(4.8), Inches(11), Inches(1.2),
                   "📌  안정적인 작업 속도 유지  +  회복과 방어 동시 진행\n\n"
                   "급하게 올리다 다시 떨어지는 악순환이 아닌,\n"
                   "한 번 올리면 유지되는 구조를 만듭니다.",
                   border_color=C_GOLD, bg_color=RGBColor(0xFF, 0xFA, 0xEB),
                   text_color=C_DARK_NAVY, font_size=14, bold=True)


def slide_09_case_ranking():
    """실적 - 순위 상승 사례"""
    slide = new_slide()
    add_bg(slide, C_DARK_NAVY)
    # 골드 라인
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Pt(4), fill_color=C_GOLD)
    add_footer(slide, 9, text_color=RGBColor(0x80, 0x90, 0xB0))

    add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                 "순위 상승 로직 운용 사례", font_size=30, bold=True, color=C_WHITE)
    add_shape(slide, Inches(1.2), Inches(1.3), Inches(2), Pt(3), fill_color=C_GOLD)

    # 사례 A - 좌측
    add_shape(slide, Inches(1.2), Inches(1.8), Inches(5.3), Inches(4.5),
              fill_color=RGBColor(0x22, 0x36, 0x5A), line_color=RGBColor(0x3A, 0x50, 0x78),
              line_width=Pt(1))
    add_text_box(slide, Inches(1.5), Inches(2.0), Inches(4.5), Inches(0.4),
                 "사례 A  |  정형외과", font_size=14, color=C_GOLD, bold=True)
    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(4.5), Inches(0.8),
                 "8위 → 2위", font_size=40, bold=True, color=C_WHITE)
    add_text_box(slide, Inches(1.5), Inches(3.3), Inches(4.5), Inches(0.4),
                 "N2  0.4816", font_size=20, bold=True, color=C_GOLD)
    add_image_placeholder(slide, Inches(1.5), Inches(3.9), Inches(4.7), Inches(2.2),
                          "사례 A 순위 그래프")

    # 사례 B - 우측
    add_shape(slide, Inches(6.9), Inches(1.8), Inches(5.3), Inches(4.5),
              fill_color=RGBColor(0x22, 0x36, 0x5A), line_color=RGBColor(0x3A, 0x50, 0x78),
              line_width=Pt(1))
    add_text_box(slide, Inches(7.2), Inches(2.0), Inches(4.5), Inches(0.4),
                 "사례 B  |  피부과", font_size=14, color=C_GOLD, bold=True)
    add_text_box(slide, Inches(7.2), Inches(2.5), Inches(4.5), Inches(0.8),
                 "1위 유지", font_size=40, bold=True, color=C_WHITE)
    add_text_box(slide, Inches(7.2), Inches(3.3), Inches(4.5), Inches(0.4),
                 "N2  0.5189", font_size=20, bold=True, color=C_GOLD)
    add_image_placeholder(slide, Inches(7.2), Inches(3.9), Inches(4.7), Inches(2.2),
                          "사례 B 순위 그래프")


def slide_10_case_recovery():
    """실적 - 3월 회복 사례"""
    slide = new_slide()
    add_bg(slide, C_WHITE)
    add_top_line(slide)
    add_footer(slide, 10)

    add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                 "3월 블로그 삭제 이후 회복 사례", font_size=28, bold=True, color=C_DARK_NAVY)
    add_shape(slide, Inches(1.2), Inches(1.3), Inches(2), Pt(3), fill_color=C_BLUE)

    cases = [
        ("사례 C", "부산피부과", "현재 14위", "N2  0.4959"),
        ("사례 D", "대구성형외과", "현재 14위", "N2  0.4148"),
        ("사례 E", "신림성형외과", "현재 3위", "N2  0.4371"),
    ]
    for i, (label, name, rank, n2) in enumerate(cases):
        cx = Inches(1.2 + i * 3.8)
        # 카드
        add_shape(slide, cx, Inches(1.7), Inches(3.4), Inches(4.8),
                  fill_color=C_LIGHT_GRAY, line_color=C_BLUE, line_width=Pt(1.5))
        add_text_box(slide, cx + Inches(0.2), Inches(1.85), Inches(3), Inches(0.35),
                     label, font_size=11, color=C_BLUE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, cx + Inches(0.2), Inches(2.2), Inches(3), Inches(0.4),
                     name, font_size=16, bold=True, color=C_DARK_NAVY,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, cx + Inches(0.2), Inches(2.65), Inches(3), Inches(0.35),
                     rank, font_size=14, color=C_DARK_TEXT, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, cx + Inches(0.2), Inches(3.0), Inches(3), Inches(0.35),
                     n2, font_size=16, bold=True, color=C_GOLD,
                     alignment=PP_ALIGN.CENTER)
        add_image_placeholder(slide, cx + Inches(0.2), Inches(3.5), Inches(3), Inches(2.8),
                              f"{label} 그래프")


def slide_11_case_salon():
    """실적 - 미용실 업종"""
    slide = new_slide()
    add_bg(slide, C_LIGHT_GRAY)
    add_top_line(slide)
    add_footer(slide, 11)

    add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.5),
                 "미용실 업종 레퍼런스", font_size=28, bold=True, color=C_DARK_NAVY)
    add_text_box(slide, Inches(1.2), Inches(1.15), Inches(10), Inches(0.35),
                 "귀사와 동일 업종 실적", font_size=15, color=C_BLUE)
    add_shape(slide, Inches(1.2), Inches(1.55), Inches(2), Pt(3), fill_color=C_BLUE)

    # 사례 F - 좌측
    add_shape(slide, Inches(1.2), Inches(1.9), Inches(5.3), Inches(1.5),
              fill_color=C_WHITE, line_color=C_BLUE, line_width=Pt(1.5))
    add_text_box(slide, Inches(1.5), Inches(2.0), Inches(2), Inches(0.35),
                 "사례 F  |  미용실", font_size=12, color=C_BLUE, bold=True)
    add_text_box(slide, Inches(1.5), Inches(2.4), Inches(4.5), Inches(0.5),
                 "15위 → 2위  (1주 13계단 ↑)", font_size=22, bold=True, color=C_DARK_NAVY)
    add_text_box(slide, Inches(1.5), Inches(2.95), Inches(4.5), Inches(0.35),
                 "N2  0.4541", font_size=16, bold=True, color=C_GOLD)

    # 사례 G - 우측
    add_shape(slide, Inches(6.9), Inches(1.9), Inches(5.3), Inches(1.5),
              fill_color=C_WHITE, line_color=C_BLUE, line_width=Pt(1.5))
    add_text_box(slide, Inches(7.2), Inches(2.0), Inches(2), Inches(0.35),
                 "사례 G  |  미용실", font_size=12, color=C_BLUE, bold=True)
    add_text_box(slide, Inches(7.2), Inches(2.4), Inches(4.5), Inches(0.5),
                 "81위 → 40위  (2주 41계단 ↑)", font_size=22, bold=True, color=C_DARK_NAVY)
    add_text_box(slide, Inches(7.2), Inches(2.95), Inches(4.5), Inches(0.35),
                 "N2  0.4460", font_size=16, bold=True, color=C_GOLD)

    # 이미지 플레이스홀더 2개
    add_image_placeholder(slide, Inches(1.2), Inches(3.7), Inches(5.3), Inches(2.2),
                          "사례 F 순위 그래프")
    add_image_placeholder(slide, Inches(6.9), Inches(3.7), Inches(5.3), Inches(2.2),
                          "사례 G 순위 그래프")

    # 하단 종합
    add_accent_box(slide, Inches(1.2), Inches(6.15), Inches(11), Inches(0.7),
                   "📊  7건 사례 종합:  순위 상승 로직 검증 완료  |  블로그 삭제 후 회복 실적 보유  |  미용실 업종 성과 입증",
                   border_color=C_DARK_NAVY, bg_color=C_WHITE,
                   text_color=C_DARK_NAVY, font_size=13, bold=True)


def slide_12_keyword_strategy():
    """키워드 전략 - 낙수효과"""
    slide = new_slide()
    add_bg(slide, C_WHITE)
    add_top_line(slide)
    add_footer(slide, 12)

    add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                 "메인 키워드 하나면 충분합니다", font_size=28, bold=True, color=C_DARK_NAVY)
    add_shape(slide, Inches(1.2), Inches(1.3), Inches(2), Pt(3), fill_color=C_BLUE)

    add_text_box(slide, Inches(1.2), Inches(1.6), Inches(11), Inches(0.8),
                 "플레이스 점수를 올리면 연관 키워드도 동반 상승합니다 (낙수효과)\n"
                 "메인 키워드 1개에 집중하면, 나머지는 자연스럽게 따라옵니다.",
                 font_size=14, color=C_DARK_TEXT, line_spacing=1.6)

    # 수치 강조 카드
    cards = [
        ("나성동미용실", "월 검색량 3,980건", "8위 → 6위"),
        ("세종미용실", "월 검색량 9,170건", "16위 → 13위"),
    ]
    for i, (kw, vol, rank) in enumerate(cards):
        cx = Inches(1.2 + i * 5.7)
        add_shape(slide, cx, Inches(2.8), Inches(5.2), Inches(1.5),
                  fill_color=C_LIGHT_GRAY, line_color=C_BLUE, line_width=Pt(1.5))
        add_text_box(slide, cx + Inches(0.3), Inches(2.95), Inches(4.5), Inches(0.35),
                     kw, font_size=16, bold=True, color=C_DARK_NAVY)
        add_text_box(slide, cx + Inches(0.3), Inches(3.3), Inches(2.5), Inches(0.3),
                     vol, font_size=11, color=RGBColor(0x64, 0x64, 0x64))
        add_text_box(slide, cx + Inches(2.8), Inches(3.15), Inches(2), Inches(0.5),
                     rank, font_size=22, bold=True, color=C_BLUE,
                     alignment=PP_ALIGN.RIGHT)

    # 이미지
    add_image_placeholder(slide, Inches(1.2), Inches(4.6), Inches(5.2), Inches(2),
                          "키워드 순위 변동 그래프")

    # 보충 설명
    add_multiline_text(slide, Inches(7), Inches(4.6), Inches(5.2), Inches(2),
                       [
                           ("참고: 연관도 낮은 키워드", 13, True, C_DARK_NAVY),
                           ("", 6),
                           ("'다정동미용실', '새롬동미용실' 등은", 12, False, C_DARK_TEXT),
                           ("플레이스명과 연관도가 낮아", 12, False, C_DARK_TEXT),
                           ("낙수효과가 제한적입니다.", 12, False, C_DARK_TEXT),
                           ("", 6),
                           ("→ 메인 키워드 집중이 효율적", 12, True, C_BLUE),
                       ])


def slide_13_safe_guarantee():
    """안전보장형 제안"""
    slide = new_slide()
    add_bg(slide, C_DARK_NAVY)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Pt(4), fill_color=C_GOLD)
    add_footer(slide, 13, text_color=RGBColor(0x80, 0x90, 0xB0))

    add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                 "제안 서비스:  안전보장형", font_size=30, bold=True, color=C_WHITE)
    add_shape(slide, Inches(1.2), Inches(1.3), Inches(2), Pt(3), fill_color=C_GOLD)

    add_text_box(slide, Inches(1.2), Inches(1.6), Inches(11), Inches(0.5),
                 "안전보장형  =  관리형의 안정성  +  순위보장형의 성과  →  최적의 균형점",
                 font_size=15, color=RGBColor(0xA0, 0xBB, 0xDD))

    # 내부 로직 테이블
    tbl = make_table(slide, Inches(1.2), Inches(2.3), Inches(11), 4, 2, Inches(0.7))
    tbl.columns[0].width = Inches(3)
    tbl.columns[1].width = Inches(8)

    set_cell_style(tbl.cell(0, 0), "전략 항목", font_size=13, bold=True,
                   color=C_WHITE, fill_color=C_GOLD)
    set_cell_style(tbl.cell(0, 1), "운영 방식", font_size=13, bold=True,
                   color=C_WHITE, fill_color=C_GOLD)

    rows_data = [
        ("키워드 세분화", "10유입 단위로 분할하여 패턴 감지 회피\n동일 키워드 반복 유입 방지"),
        ("유입경로 다양화", "매일 다른 키워드로 세팅\n블로그·카페·지도 등 다채널 분산"),
        ("매체·미션 다양화", "단일 반복 지양, 다양한 매체·미션 조합\nAI가 패턴으로 인식하지 못하는 자연스러운 시그널"),
    ]
    for i, (col1, col2) in enumerate(rows_data):
        bg = RGBColor(0x22, 0x36, 0x5A) if i % 2 == 0 else RGBColor(0x1B, 0x2F, 0x52)
        set_cell_style(tbl.cell(i + 1, 0), col1, font_size=13, bold=True,
                       color=C_GOLD, fill_color=bg)
        set_cell_style(tbl.cell(i + 1, 1), col2, font_size=12,
                       color=C_WHITE, fill_color=bg, alignment=PP_ALIGN.LEFT)

    # 핵심 메시지
    add_accent_box(slide, Inches(1.2), Inches(5.3), Inches(11), Inches(1),
                   "🛡️  안전보장형의 핵심 가치\n"
                   "순위를 올리되, 떨어지지 않게.  비용은 낮추되, 성과는 유지하게.",
                   border_color=C_GOLD, bg_color=RGBColor(0x22, 0x36, 0x5A),
                   text_color=C_GOLD, font_size=15, bold=True)


def slide_14_core_value():
    """핵심 가치 (임팩트 슬라이드)"""
    slide = new_slide()
    add_bg(slide, C_WHITE)
    add_top_line(slide)
    add_footer(slide, 14)

    # 상단 여백 후 중앙 배치
    add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(0.8),
                 "다른 업체들의 순위가 일제히 떨어질 때",
                 font_size=26, color=C_DARK_TEXT, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(1),
                 "귀사의 순위를 지킬 수 있는\n가장 확실한 방법",
                 font_size=36, bold=True, color=C_BLUE, alignment=PP_ALIGN.CENTER,
                 line_spacing=1.4)

    # 구분선
    add_shape(slide, Inches(5.5), Inches(4.7), Inches(2.3), Pt(2), fill_color=C_GOLD)

    add_text_box(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(0.5),
                 "안정적으로 유지되는 순위가 실질적인 매출을 만듭니다",
                 font_size=16, color=RGBColor(0x64, 0x64, 0x64),
                 alignment=PP_ALIGN.CENTER)


def slide_15_ai_lab():
    """AI로직연구소"""
    slide = new_slide()
    add_bg(slide, C_LIGHT_GRAY)
    add_top_line(slide)
    add_footer(slide, 15)

    add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.5),
                 "셀프마케팅 + AI로직연구소", font_size=28, bold=True, color=C_DARK_NAVY)
    add_text_box(slide, Inches(1.2), Inches(1.15), Inches(10), Inches(0.35),
                 "출시 예정  |  온다마케팅 독자 개발", font_size=14, color=C_BLUE)
    add_shape(slide, Inches(1.2), Inches(1.55), Inches(2), Pt(3), fill_color=C_BLUE)

    # 기능 테이블
    tbl = make_table(slide, Inches(1.2), Inches(1.9), Inches(7), 4, 2, Inches(0.65))
    tbl.columns[0].width = Inches(2.5)
    tbl.columns[1].width = Inches(4.5)

    set_cell_style(tbl.cell(0, 0), "기능", font_size=13, bold=True,
                   color=C_WHITE, fill_color=C_DARK_NAVY)
    set_cell_style(tbl.cell(0, 1), "상세 내용", font_size=13, bold=True,
                   color=C_WHITE, fill_color=C_DARK_NAVY)

    features = [
        ("실시간 로직 모니터링", "3.5일 간격으로 알고리즘 변동 감지\n선제적 대응 가능"),
        ("맞춤 매체 추천", "안전성 + 효율성 최적 매체 자동 추천\n리스크 최소화"),
        ("자동 등록 및 전환", "충전 시 자동 처리\n수동 작업 최소화"),
    ]
    for i, (feat, desc) in enumerate(features):
        bg = C_TABLE_ALT1 if i % 2 == 0 else C_TABLE_ALT2
        set_cell_style(tbl.cell(i + 1, 0), feat, font_size=12, bold=True,
                       color=C_DARK_NAVY, fill_color=bg)
        set_cell_style(tbl.cell(i + 1, 1), desc, font_size=11,
                       color=C_DARK_TEXT, fill_color=bg, alignment=PP_ALIGN.LEFT)

    # 이미지 플레이스홀더
    add_image_placeholder(slide, Inches(8.7), Inches(1.9), Inches(3.5), Inches(4.3),
                          "AI로직연구소 솔루션 화면")


def slide_16_pricing():
    """견적"""
    slide = new_slide()
    add_bg(slide, C_WHITE)
    add_top_line(slide)
    add_footer(slide, 16)

    add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                 "견적 안내", font_size=30, bold=True, color=C_DARK_NAVY)
    add_shape(slide, Inches(1.2), Inches(1.3), Inches(2), Pt(3), fill_color=C_BLUE)

    # 견적 테이블
    tbl = make_table(slide, Inches(0.8), Inches(1.7), Inches(11.7), 6, 5, Inches(0.6))
    col_widths = [Inches(2.2), Inches(1.5), Inches(2.5), Inches(2.7), Inches(2.8)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    headers = ["지점", "현재순위", "보장형 (첫달)", "보장형 (2개월~)", "안전보장형 (권장)"]
    for j, h in enumerate(headers):
        bg = C_GOLD if j == 4 else C_DARK_NAVY
        set_cell_style(tbl.cell(0, j), h, font_size=12, bold=True,
                       color=C_WHITE, fill_color=bg)

    data = [
        ["마크원애비뉴점", "13위", "1,800,000원", "1,200,000원", "900,000원"],
        ["나성1호점", "15위", "1,900,000원", "1,200,000원", "900,000원"],
        ["나성2호점", "20위", "2,100,000원", "1,200,000원", "900,000원"],
        ["반곡점", "21위", "2,200,000원", "1,200,000원", "900,000원"],
        ["4지점 합계", "—", "8,000,000원", "4,800,000원", "3,600,000원"],
    ]
    for r, row in enumerate(data):
        is_total = r == 4
        for c, val in enumerate(row):
            if is_total:
                bg = C_GOLD_LIGHT if c == 4 else RGBColor(0xE8, 0xEB, 0xF0)
                bld = True
            elif c == 4:
                bg = RGBColor(0xFF, 0xFB, 0xEB)
                bld = True
            else:
                bg = C_TABLE_ALT1 if r % 2 == 0 else C_TABLE_ALT2
                bld = False
            clr = C_DARK_NAVY if is_total or c == 4 else C_DARK_TEXT
            set_cell_style(tbl.cell(r + 1, c), val, font_size=12, bold=bld,
                           color=clr, fill_color=bg)

    # 참고
    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.8),
                 "※ 4개 지점 동시 진행 시 지점당 월 90만원 균일\n"
                 "※ 안전보장형은 계약기간 없음, 1주 단위 리포팅 제공",
                 font_size=12, color=RGBColor(0x64, 0x64, 0x64))

    # 강조
    add_accent_box(slide, Inches(0.8), Inches(6.15), Inches(11.7), Inches(0.6),
                   "💡  안전보장형 4지점 합계 월 360만원  —  보장형 대비 55% 절감 + 하락 리스크 최소화",
                   border_color=C_GOLD, bg_color=C_GOLD_LIGHT,
                   text_color=C_DARK_NAVY, font_size=14, bold=True)


def slide_17_comparison():
    """운영 조건 비교"""
    slide = new_slide()
    add_bg(slide, C_LIGHT_GRAY)
    add_top_line(slide)
    add_footer(slide, 17)

    add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                 "운영 조건 비교", font_size=28, bold=True, color=C_DARK_NAVY)
    add_shape(slide, Inches(1.2), Inches(1.3), Inches(2), Pt(3), fill_color=C_BLUE)

    # 비교 테이블
    tbl = make_table(slide, Inches(1.2), Inches(1.7), Inches(11), 8, 3, Inches(0.55))
    tbl.columns[0].width = Inches(2.5)
    tbl.columns[1].width = Inches(4.25)
    tbl.columns[2].width = Inches(4.25)

    headers = ["항목", "보장형", "안전보장형 (권장)"]
    header_colors = [C_DARK_NAVY, C_DARK_NAVY, C_BLUE]
    for j, h in enumerate(headers):
        set_cell_style(tbl.cell(0, j), h, font_size=13, bold=True,
                       color=C_WHITE, fill_color=header_colors[j])

    rows_data = [
        ("작업방식", "공격적 트래픽 투입", "내부 로직 중심 안정적 관리"),
        ("상승속도", "빠름", "점진적 / 안정적"),
        ("하락 리스크", "높음", "낮음"),
        ("결제방식", "노출 후 결제", "월 관리비"),
        ("계약기간", "없음", "없음"),
        ("리포팅", "1주일", "1주일 (대시보드 상시)"),
        ("월비용 (4지점)", "첫달 800만 / 유지 480만", "360만원 (균일)"),
    ]
    for r, (item, bj, saf) in enumerate(rows_data):
        bg = C_TABLE_ALT1 if r % 2 == 0 else C_TABLE_ALT2
        set_cell_style(tbl.cell(r + 1, 0), item, font_size=12, bold=True,
                       color=C_DARK_NAVY, fill_color=bg)
        # 보장형: 리스크 높음은 빨간색
        bj_color = C_RED if "높음" in bj else C_DARK_TEXT
        set_cell_style(tbl.cell(r + 1, 1), bj, font_size=12,
                       color=bj_color, fill_color=bg)
        # 안전보장형: 핵심 강조
        saf_color = C_GREEN if "낮음" in saf else C_BLUE if "360" in saf else C_DARK_NAVY
        saf_bold = "360" in saf or "낮음" in saf
        set_cell_style(tbl.cell(r + 1, 2), saf, font_size=12, bold=saf_bold,
                       color=saf_color, fill_color=RGBColor(0xEF, 0xF6, 0xFF))

    # 하단 권장
    add_accent_box(slide, Inches(1.2), Inches(6.1), Inches(11), Inches(0.7),
                   "온다마케팅은 안전보장형을 권장드립니다.\n"
                   "급격한 상승보다 안정적인 유지가 실질 매출에 기여합니다.",
                   border_color=C_BLUE, bg_color=RGBColor(0xEF, 0xF6, 0xFF),
                   text_color=C_DARK_NAVY, font_size=14, bold=True)


def slide_18_closing():
    """마무리"""
    slide = new_slide()
    add_bg(slide, C_DARK_NAVY)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Pt(4), fill_color=C_GOLD)

    # 중앙 감사 메시지
    add_text_box(slide, Inches(1), Inches(2.3), Inches(11.3), Inches(1),
                 "감사합니다", font_size=44, bold=True, color=C_WHITE,
                 alignment=PP_ALIGN.CENTER)

    # 구분선
    add_shape(slide, Inches(5.5), Inches(3.5), Inches(2.3), Pt(2), fill_color=C_GOLD)

    add_text_box(slide, Inches(1), Inches(3.9), Inches(11.3), Inches(0.7),
                 "어떤 방식을 선택하시더라도\n최선의 결과를 만들겠습니다",
                 font_size=20, color=C_GOLD, alignment=PP_ALIGN.CENTER,
                 line_spacing=1.5)

    # 하단 연락처
    add_shape(slide, Inches(3.5), Inches(5.5), Inches(6.3), Pt(1),
              fill_color=RGBColor(0x3A, 0x50, 0x78))
    add_text_box(slide, Inches(1), Inches(5.8), Inches(11.3), Inches(0.5),
                 "온다마케팅  |  연락처: [담당자 연락처 기입]",
                 font_size=13, color=RGBColor(0x80, 0x90, 0xB0),
                 alignment=PP_ALIGN.CENTER)


# ── 메인 실행 ──

def main():
    print("🔨 마드모아젤헤어 제안서 PPT 생성 시작...")

    slide_01_cover()
    print("  ✅ 슬라이드 1/18: 표지")
    slide_02_greeting()
    print("  ✅ 슬라이드 2/18: 인사말")
    slide_03_rule_change()
    print("  ✅ 슬라이드 3/18: 2026년 규칙 변화")
    slide_04_defense_system()
    print("  ✅ 슬라이드 4/18: 대응 체계")
    slide_05_diagnosis()
    print("  ✅ 슬라이드 5/18: 현황 진단")
    slide_06_qa_n2()
    print("  ✅ 슬라이드 6/18: Q&A N2 하락 원인")
    slide_07_qa_blog()
    print("  ✅ 슬라이드 7/18: Q&A 블로그 삭제")
    slide_08_qa_recovery()
    print("  ✅ 슬라이드 8/18: Q&A 회복 플랜")
    slide_09_case_ranking()
    print("  ✅ 슬라이드 9/18: 순위 상승 사례")
    slide_10_case_recovery()
    print("  ✅ 슬라이드 10/18: 3월 회복 사례")
    slide_11_case_salon()
    print("  ✅ 슬라이드 11/18: 미용실 업종")
    slide_12_keyword_strategy()
    print("  ✅ 슬라이드 12/18: 키워드 전략")
    slide_13_safe_guarantee()
    print("  ✅ 슬라이드 13/18: 안전보장형 제안")
    slide_14_core_value()
    print("  ✅ 슬라이드 14/18: 핵심 가치")
    slide_15_ai_lab()
    print("  ✅ 슬라이드 15/18: AI로직연구소")
    slide_16_pricing()
    print("  ✅ 슬라이드 16/18: 견적")
    slide_17_comparison()
    print("  ✅ 슬라이드 17/18: 운영 조건 비교")
    slide_18_closing()
    print("  ✅ 슬라이드 18/18: 마무리")

    output_path = "/home/onda/projects/onda-hompage/마드모아젤헤어_제안서.pptx"
    prs.save(output_path)
    print(f"\n✅ PPT 생성 완료: {output_path}")
    print(f"   슬라이드 수: {len(prs.slides)}장")


if __name__ == "__main__":
    main()
