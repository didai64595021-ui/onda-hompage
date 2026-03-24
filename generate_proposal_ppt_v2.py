#!/usr/bin/env python3
"""마드모아젤헤어 제안서 PPT 생성 - 고도화 버전"""

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === 그리드 시스템 ===
SLIDE_W = Cm(33.867)
SLIDE_H = Cm(19.05)
ML = Cm(2.5)
MR = Cm(2.5)
CW = Cm(28.867)
TITLE_TOP = Cm(1.8)
TITLE_H = Cm(1.3)
ACCENT_TOP = Cm(3.3)
BODY_TOP = Cm(4.0)
FOOTER_TOP = Cm(17.8)
LOGO_W = Cm(4.0)
LOGO_H = Cm(1.2)
LOGO_LEFT = Cm(33.867 - 2.5 - 4.0)
LOGO_TOP = Cm(0.5)

# === 색상 ===
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x25, 0x63, 0xEB)
GOLD = RGBColor(0xD4, 0xA8, 0x53)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
TABLE_ALT = RGBColor(0xF0, 0xF4, 0xFF)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
PH_BG = RGBColor(0xF3, 0xF4, 0xF6)
PH_BORDER = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_TEXT = RGBColor(0x6B, 0x72, 0x80)

LOGO_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'onda_logo.jpg')

# === 유틸 함수 ===
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_top_line(slide, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), SLIDE_W, Cm(0.12))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_footer(slide, page, total=18):
    tf = slide.shapes.add_textbox(Cm(28), FOOTER_TOP, Cm(5), Cm(0.7)).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"온다마케팅 | {page}/{total}"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY_TEXT
    p.alignment = PP_ALIGN.RIGHT

def add_footer_dark(slide, page, total=18):
    tf = slide.shapes.add_textbox(Cm(28), FOOTER_TOP, Cm(5), Cm(0.7)).text_frame
    p = tf.paragraphs[0]
    p.text = f"온다마케팅 | {page}/{total}"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0xAA)
    p.alignment = PP_ALIGN.RIGHT

def add_logo(slide):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, LOGO_LEFT, LOGO_TOP, LOGO_W, LOGO_H)

def add_title(slide, text, color=DARK_TEXT, size=Pt(28)):
    tf = slide.shapes.add_textbox(ML, TITLE_TOP, CW, TITLE_H).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "맑은 고딕"

def add_accent(slide, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, ML, ACCENT_TOP, Cm(4.0), Cm(0.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def tb(slide, left, top, w, h, text, size=Pt(12), bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    tf = slide.shapes.add_textbox(left, top, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def multi_tb(slide, left, top, w, h, lines, default_size=Pt(12), default_color=DARK_TEXT):
    """여러 줄 텍스트박스. lines = [(text, size, bold, color), ...]"""
    tf = slide.shapes.add_textbox(left, top, w, h).text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        text, size, bold, color = line[0], line[1] if len(line)>1 else default_size, line[2] if len(line)>2 else False, line[3] if len(line)>3 else default_color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "맑은 고딕"
        p.space_after = Pt(4)
    return tf

def highlight_box(slide, left, top, w, h, text, bg, border, text_color=DARK_TEXT, size=Pt(12)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = border
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.4)
    tf.margin_right = Cm(0.4)
    tf.margin_top = Cm(0.2)
    tf.margin_bottom = Cm(0.2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = text_color
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.LEFT
    return shape

def placeholder_img(slide, left, top, w, h, label="이미지 자리"):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PH_BG
    shape.line.color.rgb = PH_BORDER
    shape.line.width = Pt(1.5)
    shape.line.dash_style = 2  # dash
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"[{label}]"
    p.font.size = Pt(11)
    p.font.color.rgb = PH_BORDER
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def qa_icon(slide, left, top, letter, bg_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Cm(1.5), Cm(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Cm(0)
    tf.margin_right = Cm(0)
    tf.margin_top = Cm(0.1)
    p = tf.paragraphs[0]
    p.text = letter
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER

def styled_table(slide, left, top, w, rows, cols, data, col_widths=None):
    """스타일 테이블 생성"""
    row_h = Cm(0.9)
    table_h = row_h * rows
    shape = slide.shapes.add_table(rows, cols, left, top, w, table_h)
    tbl = shape.table
    
    if col_widths:
        for i, cw_val in enumerate(col_widths):
            tbl.columns[i].width = cw_val
    
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ""
            cell.margin_left = Cm(0.2)
            cell.margin_right = Cm(0.2)
            cell.margin_top = Cm(0.1)
            cell.margin_bottom = Cm(0.1)
            
            for p in cell.text_frame.paragraphs:
                p.font.name = "맑은 고딕"
                p.font.size = Pt(10)
                p.alignment = PP_ALIGN.CENTER
            
            # 헤더
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = WHITE
                    p.font.bold = True
                    p.font.size = Pt(10)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT if r % 2 == 0 else WHITE
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = DARK_TEXT
    
    return shape

def card_box(slide, left, top, w, h, top_color=BLUE):
    """카드 박스 (상단 컬러 바 포함)"""
    # 메인 박스
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    shape.line.width = Pt(1)
    # 상단 바
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Cm(0.25))
    bar.fill.solid()
    bar.fill.fore_color.rgb = top_color
    bar.line.fill.background()
    return shape

# === 슬라이드 생성 ===
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]  # blank layout

# ====== SLIDE 1: 표지 ======
s = prs.slides.add_slide(blank)
set_bg(s, NAVY)
add_top_line(s, GOLD)
# 로고 중앙
if os.path.exists(LOGO_PATH):
    s.shapes.add_picture(LOGO_PATH, Cm(12), Cm(2.5), Cm(10), Cm(3))
tb(s, ML, Cm(6.5), CW, Cm(0.8), "네이버 종합광고 실행사", Pt(13), False, RGBColor(0xBB,0xBB,0xCC), PP_ALIGN.CENTER)
tb(s, ML, Cm(7.8), CW, Cm(2.0), "마드모아젤헤어", Pt(42), True, WHITE, PP_ALIGN.CENTER)
tb(s, ML, Cm(10.2), CW, Cm(1.2), "플레이스 순위 회복 및 방어 전략 제안서", Pt(20), False, GOLD, PP_ALIGN.CENTER)
tb(s, ML, Cm(12.0), CW, Cm(0.8), "세종시 4개 지점 통합 진단 | 안전보장형 운영 제안", Pt(13), False, RGBColor(0x99,0x99,0xAA), PP_ALIGN.CENTER)
tb(s, Cm(28), Cm(17.5), Cm(5), Cm(0.6), "2026.03", Pt(10), False, RGBColor(0x88,0x88,0x99), PP_ALIGN.RIGHT)

# ====== SLIDE 2: 인사말 ======
s = prs.slides.add_slide(blank)
add_top_line(s)
add_logo(s)
add_title(s, "인사말")
add_accent(s)
tb(s, ML, BODY_TOP, CW, Cm(4.0),
   "안녕하세요. 마드모아젤헤어 마케팅 미팅에 참여하게 되어 감사드립니다.\n\n"
   "저희는 네이버 플레이스, 인스타그램 등 온라인 채널을 통합 운영하는\n"
   "종합광고 실행사 온다마케팅입니다.", Pt(13), False, DARK_TEXT)
highlight_box(s, ML, Cm(9.5), CW, Cm(2.5),
    "본 제안서는 귀사의 현재 플레이스 현황에 대한 진단과 함께,\n"
    "2026년 변화된 알고리즘 환경에 맞는 대응 전략을 담고 있습니다.",
    RGBColor(0xEF,0xF6,0xFF), BLUE, DARK_TEXT, Pt(12))
add_footer(s, 2)

# ====== SLIDE 3: 규칙 변화 ======
s = prs.slides.add_slide(blank)
set_bg(s, LIGHT_BG)
add_top_line(s)
add_logo(s)
add_title(s, "2026년, 플레이스 마케팅의 규칙이 바뀌었습니다")
add_accent(s)

# 좌측: 과거
col_w = Cm(13.5)
gap = Cm(1.867)
left1 = ML
left2 = ML + col_w + gap

# 과거 박스
shape = slide_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left1, Cm(4.2), col_w, Cm(9.0))
shape.fill.solid()
shape.fill.fore_color.rgb = WHITE
shape.line.color.rgb = RGBColor(0xDD,0xDD,0xDD)
shape.line.width = Pt(1)
tb(s, left1+Cm(0.5), Cm(4.5), col_w-Cm(1), Cm(1.0), "과거 (2025년 이전)", Pt(16), True, GRAY_TEXT, PP_ALIGN.CENTER)
past_items = [
    "• 트래픽 양 = 순위",
    "• 3~6개월 단위 로직 변경",
    "• 배포형 블로그로 순위 상승 가능",
    "• 공격적 작업 = 빠른 성과"
]
multi_tb(s, left1+Cm(0.8), Cm(5.8), col_w-Cm(1.6), Cm(6.0),
    [(t, Pt(12), False, DARK_TEXT) for t in past_items])

# 현재 박스
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left2, Cm(4.2), col_w, Cm(9.0))
shape.fill.solid()
shape.fill.fore_color.rgb = WHITE
shape.line.color.rgb = BLUE
shape.line.width = Pt(2)
tb(s, left2+Cm(0.5), Cm(4.5), col_w-Cm(1), Cm(1.0), "현재 (2026년 이후)", Pt(16), True, BLUE, PP_ALIGN.CENTER)
now_items = [
    "• 트래픽 패턴 감지 → 제재/패널티",
    "• 1~2주 단위 세부 로직 변동",
    "• 배포형 블로그 대량 삭제 + 순위 급락",
    "• 공격적 작업 = 공격적 하락 수반"
]
multi_tb(s, left2+Cm(0.8), Cm(5.8), col_w-Cm(1.6), Cm(6.0),
    [(t, Pt(12), False, DARK_TEXT) for t in now_items])

# 화살표
tb(s, Cm(16.2), Cm(7.5), Cm(1.5), Cm(1.5), "→", Pt(36), True, BLUE, PP_ALIGN.CENTER)

# 하단 강조
highlight_box(s, ML, Cm(14.0), CW, Cm(1.8),
    "핵심: 순위를 '올리는 것'보다, '지키는 것'이 더 어렵고 더 중요한 시대입니다.",
    RGBColor(0xFE,0xF9,0xEF), GOLD, DARK_TEXT, Pt(13))
add_footer(s, 3)

# ====== SLIDE 4: 대응 체계 ======
s = prs.slides.add_slide(blank)
add_top_line(s)
add_logo(s)
add_title(s, "데이터 기반 순위 방어 체계")
add_accent(s)

card_w = Cm(8.79)  # (CW - 2*gap) / 3
card_gap = Cm(0.65)
cards_data = [
    ("🔬", "셀프마케팅 +\nAI로직연구소 구축", "상세 내용 9장 참조"),
    ("🛡️", "순위 방어 +\n상승 로직 적용", "구체화된 로직 기반 운영"),
    ("📈", "안전하고 꾸준한\n상승 추구", "과도한 트래픽 대신 안정적 방어"),
]
for i, (icon, title, desc) in enumerate(cards_data):
    x = ML + i * (card_w + card_gap)
    card_box(s, x, Cm(4.5), card_w, Cm(8.5), BLUE)
    tb(s, x, Cm(5.2), card_w, Cm(2.0), icon, Pt(40), False, BLUE, PP_ALIGN.CENTER)
    tb(s, x+Cm(0.5), Cm(7.5), card_w-Cm(1), Cm(2.0), title, Pt(14), True, DARK_TEXT, PP_ALIGN.CENTER)
    tb(s, x+Cm(0.5), Cm(10.0), card_w-Cm(1), Cm(2.0), desc, Pt(11), False, GRAY_TEXT, PP_ALIGN.CENTER)

tb(s, ML, Cm(14.0), CW, Cm(1.5),
   "단순 트래픽 유입 대행에서 벗어나, 실시간 데이터 분석 기반의 순위 방어 체계를 구축하고 있습니다.",
   Pt(12), False, DARK_TEXT, PP_ALIGN.CENTER)
add_footer(s, 4)

# ====== SLIDE 5: 현황 진단 ======
s = prs.slides.add_slide(blank)
add_top_line(s)
add_logo(s)
add_title(s, "마드모아젤헤어 4개 지점 현황 진단")
add_accent(s)

data5 = [
    ["지점", "순위 변동", "N2 변동", "블로그 삭제", "저품질"],
    ["마크원애비뉴점", "하락", "-0.08 이상", "해당", "없음"],
    ["나성1호점", "1위→14위", "-0.08 이상", "해당", "없음"],
    ["나성2호점", "14위→21위", "-0.10", "해당", "없음"],
    ["반곡점", "하락", "-0.08 이상", "해당", "없음"],
]
tbl_shape = styled_table(s, ML, BODY_TOP, CW, 5, 5, data5)

tb(s, ML, Cm(8.8), CW, Cm(2.0),
   "공통: N1 변동 없음. N2 전 지점 동시 급락.\n블로그 리뷰 약 300~400건 삭제/누락. 방문자 리뷰 증가에도 순위 방어 실패.",
   Pt(11), False, DARK_TEXT)

highlight_box(s, ML, Cm(11.5), CW, Cm(2.0),
    "✅ 종합 판단: 4개 지점 모두 저품질 상태가 아닙니다.\n회복 가능한 상태이며, 적절한 대응 시 순위 정상화가 가능합니다.",
    RGBColor(0xEC,0xFD,0xF5), GREEN, DARK_TEXT, Pt(12))
add_footer(s, 5)

# ====== SLIDE 6: Q&A #1 ======
s = prs.slides.add_slide(blank)
set_bg(s, LIGHT_BG)
add_top_line(s)
add_logo(s)

qa_icon(s, ML, Cm(1.8), "Q", BLUE)
tb(s, ML+Cm(2.0), Cm(1.9), CW-Cm(2.0), Cm(1.3),
   "전 지점 N2 동시 하락, 원인이 무엇입니까?", Pt(18), True, DARK_TEXT)

qa_icon(s, ML, Cm(4.0), "A", NAVY)
tb(s, ML+Cm(2.0), Cm(4.1), CW-Cm(2.0), Cm(5.0),
   "가장 직접적인 원인은 3월 초 배포형 블로그 및 비실명 계정의\n"
   "대대적인 삭제 조치입니다.\n\n"
   "네이버는 각 플레이스별로 삭제량과 키워드에 따라 회복 기간을\n"
   "차등 적용하고 있으며, 저희는 자사 관리 업체 데이터를 토대로\n"
   "이 회복 패턴을 정확히 파악하고 있습니다.",
   Pt(12), False, DARK_TEXT)

# 핵심 포인트 3개
points = [
    ("✓  N1(SEO) 점수 변동 없음", GREEN),
    ("✓  4개 지점 모두 저품질 해당 없음", GREEN),
    ("✓  정상적인 회복 가능", GREEN),
]
for i, (txt, clr) in enumerate(points):
    highlight_box(s, ML+Cm(2.0), Cm(10.5)+Cm(i*1.6), Cm(15), Cm(1.2),
        txt, RGBColor(0xEC,0xFD,0xF5), GREEN, DARK_TEXT, Pt(12))
add_footer(s, 6)

# ====== SLIDE 7: Q&A #2 ======
s = prs.slides.add_slide(blank)
add_top_line(s)
add_logo(s)

qa_icon(s, ML, Cm(1.8), "Q", BLUE)
tb(s, ML+Cm(2.0), Cm(1.9), CW-Cm(2.0), Cm(1.3),
   "블로그 대량 삭제, 어떻게 상쇄합니까?", Pt(18), True, DARK_TEXT)

qa_icon(s, ML, Cm(4.0), "A", NAVY)
tb(s, ML+Cm(2.0), Cm(4.1), CW-Cm(2.0), Cm(6.0),
   "단순 재배포는 오히려 독입니다.\n\n"
   "현재는 2주 단위, 빠르면 1주 단위로 내부 가중치가 변동되고 있습니다.\n"
   "CLOVA AI 도입 이후 유입 경로, fingerprint 유사성, 패턴이 있는 방식의\n"
   "작업은 효율이 급락하고 순위 하락으로 이어지고 있습니다.",
   Pt(12), False, DARK_TEXT)

highlight_box(s, ML+Cm(2.0), Cm(11.0), CW-Cm(2.0), Cm(2.0),
    "이는 6장에 첨부해 드린 당사 관리 업체의 실제 데이터를 통해서도\n확인하실 수 있습니다.",
    RGBColor(0xEF,0xF6,0xFF), BLUE, DARK_TEXT, Pt(12))
add_footer(s, 7)

# ====== SLIDE 8: Q&A #3 ======
s = prs.slides.add_slide(blank)
set_bg(s, LIGHT_BG)
add_top_line(s)
add_logo(s)

qa_icon(s, ML, Cm(1.8), "Q", BLUE)
tb(s, ML+Cm(2.0), Cm(1.9), CW-Cm(2.0), Cm(1.3),
   "급락 지점의 1단계 회복 플랜은?", Pt(18), True, DARK_TEXT)

qa_icon(s, ML, Cm(4.0), "A", NAVY)
tb(s, ML+Cm(2.0), Cm(4.1), CW-Cm(2.0), Cm(5.0),
   "패턴이 드러날 수밖에 없는 단순 트래픽 작업은 동일한 리스크를 반복합니다.\n"
   "그러나 경쟁사 전부가 트래픽을 넣고 있는 상황에서 트래픽 없는 순위 상승 역시 불가능합니다.\n\n"
   "핵심 접근: 패턴화되지 않는 작업, AI가 감지할 수 없는 방식으로\n"
   "트래픽을 설계합니다.",
   Pt(12), False, DARK_TEXT)

highlight_box(s, ML, Cm(10.5), CW, Cm(2.5),
    "저희가 약속드리는 것:\n"
    "지나치게 공격적이지도, 지나치게 느리지도 않은 안정적인 작업 속도를 유지하면서,\n"
    "회복과 방어를 동시에 진행하겠습니다.",
    RGBColor(0xFE,0xF9,0xEF), GOLD, DARK_TEXT, Pt(12))
add_footer(s, 8)

# ====== SLIDE 9: 순위 상승 사례 ======
s = prs.slides.add_slide(blank)
set_bg(s, NAVY)
add_top_line(s, GOLD)

tb(s, ML, Cm(1.5), CW, Cm(1.3),
   "순위 상승 로직 운용 사례", Pt(28), True, WHITE, PP_ALIGN.LEFT)
tb(s, ML, Cm(3.0), CW, Cm(0.8),
   "안정적으로 순위 상승을 이끌어내는 로직을 운용 중입니다.", Pt(14), False, GOLD)

col_w9 = Cm(13.5)
gap9 = Cm(1.867)
# 사례 A
tb(s, ML, Cm(4.5), col_w9, Cm(0.7), "사례 A. 정형외과 키워드", Pt(13), False, RGBColor(0xBB,0xBB,0xCC))
tb(s, ML, Cm(5.3), col_w9, Cm(1.5), "8위 → 2위", Pt(36), True, GOLD, PP_ALIGN.LEFT)
tb(s, ML, Cm(7.0), col_w9, Cm(0.7), "N2 0.4816 (상승세)", Pt(12), False, WHITE)
placeholder_img(s, ML, Cm(8.0), col_w9, Cm(7.0), "이미지 자리: 사례 A 그래프")

# 사례 B
x2 = ML + col_w9 + gap9
tb(s, x2, Cm(4.5), col_w9, Cm(0.7), "사례 B. 피부과 키워드", Pt(13), False, RGBColor(0xBB,0xBB,0xCC))
tb(s, x2, Cm(5.3), col_w9, Cm(1.5), "1위 유지", Pt(36), True, GOLD, PP_ALIGN.LEFT)
tb(s, x2, Cm(7.0), col_w9, Cm(0.7), "N2 0.5189 | 경쟁 키워드 내 N2 1위", Pt(12), False, WHITE)
placeholder_img(s, x2, Cm(8.0), col_w9, Cm(7.0), "이미지 자리: 사례 B 그래프")
add_footer_dark(s, 9)

# ====== SLIDE 10: 회복 사례 ======
s = prs.slides.add_slide(blank)
add_top_line(s)
add_logo(s)
add_title(s, "3월 블로그 삭제 이후 회복 사례")
add_accent(s)

card_w10 = Cm(8.79)
card_gap10 = Cm(0.65)
cases = [
    ("사례 C", "부산피부과", "14위", "N2 0.4959 (회복 중)", "블로그 6,121건\n리뷰 2,588건"),
    ("사례 D", "대구성형외과", "14위", "N2 0.4148 (회복 중)", "블로그 330건\n리뷰 726건"),
    ("사례 E", "신림성형외과", "3위", "N2 0.4371 (결제 진행 중)", "블로그 627건\n리뷰 363건"),
]
for i, (label, kw, rank, n2, detail) in enumerate(cases):
    x = ML + i * (card_w10 + card_gap10)
    card_box(s, x, Cm(4.2), card_w10, Cm(12.0), BLUE)
    tb(s, x+Cm(0.3), Cm(4.8), card_w10-Cm(0.6), Cm(0.6), label, Pt(10), False, GRAY_TEXT, PP_ALIGN.LEFT)
    tb(s, x+Cm(0.3), Cm(5.3), card_w10-Cm(0.6), Cm(0.7), kw, Pt(14), True, DARK_TEXT, PP_ALIGN.LEFT)
    tb(s, x+Cm(0.3), Cm(6.2), card_w10-Cm(0.6), Cm(1.2), rank, Pt(28), True, BLUE, PP_ALIGN.CENTER)
    tb(s, x+Cm(0.3), Cm(7.6), card_w10-Cm(0.6), Cm(0.6), n2, Pt(11), False, DARK_TEXT, PP_ALIGN.CENTER)
    tb(s, x+Cm(0.3), Cm(8.4), card_w10-Cm(0.6), Cm(1.0), detail, Pt(10), False, GRAY_TEXT, PP_ALIGN.CENTER)
    placeholder_img(s, x+Cm(0.3), Cm(9.8), card_w10-Cm(0.6), Cm(5.8), f"이미지 자리: {label} 그래프")
add_footer(s, 10)

# ====== SLIDE 11: 미용실 업종 ======
s = prs.slides.add_slide(blank)
set_bg(s, LIGHT_BG)
add_top_line(s)
add_logo(s)
add_title(s, "미용실 업종 레퍼런스")
tb(s, ML, Cm(3.0), CW, Cm(0.7), "귀사와 동일한 업종에서의 실제 운영 실적입니다.", Pt(13), False, GRAY_TEXT)

col_w11 = Cm(13.5)
gap11 = Cm(1.867)
# 사례 F
card_box(s, ML, Cm(4.2), col_w11, Cm(5.0), BLUE)
tb(s, ML+Cm(0.5), Cm(4.8), col_w11-Cm(1), Cm(0.6), "사례 F. 미용실 키워드", Pt(12), False, GRAY_TEXT)
tb(s, ML+Cm(0.5), Cm(5.5), col_w11-Cm(1), Cm(1.5), "15위 → 2위", Pt(32), True, BLUE)
tb(s, ML+Cm(0.5), Cm(7.2), col_w11-Cm(1), Cm(1.0), "1주 만에 13계단 상승 | N2 0.4541\n블로그 402건 | 리뷰 606건", Pt(11), False, DARK_TEXT)

# 사례 G
x2_11 = ML + col_w11 + gap11
card_box(s, x2_11, Cm(4.2), col_w11, Cm(5.0), BLUE)
tb(s, x2_11+Cm(0.5), Cm(4.8), col_w11-Cm(1), Cm(0.6), "사례 G. 미용실 키워드", Pt(12), False, GRAY_TEXT)
tb(s, x2_11+Cm(0.5), Cm(5.5), col_w11-Cm(1), Cm(1.5), "81위 → 40위", Pt(32), True, BLUE)
tb(s, x2_11+Cm(0.5), Cm(7.2), col_w11-Cm(1), Cm(1.0), "2주간 41계단 상승 | N2 0.4460\n블로그 257건 | 리뷰 13,780건", Pt(11), False, DARK_TEXT)

# 플레이스홀더
placeholder_img(s, ML, Cm(10.0), col_w11, Cm(5.0), "이미지 자리: 사례 F 그래프")
placeholder_img(s, x2_11, Cm(10.0), col_w11, Cm(5.0), "이미지 자리: 사례 G 그래프")

highlight_box(s, ML, Cm(15.8), CW, Cm(1.5),
    "7건 사례 종합: 순위 상승 로직 보유(A,B,F,G) + 회복 실적(C,D,E) + 미용실 성과(F,G)",
    RGBColor(0xEF,0xF6,0xFF), BLUE, DARK_TEXT, Pt(11))
add_footer(s, 11)

# ====== SLIDE 12: 낙수효과 ======
s = prs.slides.add_slide(blank)
add_top_line(s)
add_logo(s)
add_title(s, "메인 키워드 하나면 충분합니다")
add_accent(s)

tb(s, ML, BODY_TOP, CW, Cm(2.0),
   "플레이스 마케팅은 '키워드'를 올리는 작업이 아니라\n"
   "플레이스 자체의 점수를 올리는 작업입니다.\n"
   "점수가 올라가면 연관된 모든 키워드에서 순위가 동반 상승합니다(낙수효과).",
   Pt(12), False, DARK_TEXT)

# 수치 카드 2개
kw_w = Cm(13.5)
kw_gap = Cm(1.867)
for i, (kw, vol, change, clr) in enumerate([
    ("나성동미용실", "3,980건", "8위 → 6위", GREEN),
    ("세종미용실", "9,170건", "16위 → 13위", GREEN),
]):
    x = ML + i * (kw_w + kw_gap)
    card_box(s, x, Cm(7.0), kw_w, Cm(3.0), clr)
    tb(s, x+Cm(0.5), Cm(7.5), kw_w-Cm(1), Cm(0.6), f"{kw} ({vol})", Pt(12), False, DARK_TEXT, PP_ALIGN.CENTER)
    tb(s, x+Cm(0.5), Cm(8.3), kw_w-Cm(1), Cm(1.2), change, Pt(26), True, clr, PP_ALIGN.CENTER)

placeholder_img(s, ML, Cm(10.8), Cm(18), Cm(4.5), "이미지 자리: 마크원애비뉴점 낙수효과 데이터")

tb(s, ML, Cm(15.8), CW, Cm(1.5),
   "'다정동미용실', '새롬동미용실' 등에서 순위가 낮은 것은 소재지 연관도 차이이며 정상적인 로직입니다.",
   Pt(11), False, GRAY_TEXT)
add_footer(s, 12)

# ====== SLIDE 13: 안전보장형 ======
s = prs.slides.add_slide(blank)
set_bg(s, NAVY)
add_top_line(s, GOLD)

tb(s, ML, Cm(1.5), CW, Cm(1.3), "제안 서비스: 안전보장형", Pt(28), True, WHITE)
tb(s, ML, Cm(3.0), CW, Cm(0.8), "관리형과 순위보장형의 중간 형태", Pt(14), False, GOLD)

# 안전보장형 설명
tb(s, ML, Cm(4.5), CW, Cm(2.5),
   "단순보장형도 가능합니다. 6장 사례처럼, 저희는 가파른 순위 상승 로직을 보유하고\n"
   "실제 운영 중입니다. 다만 보장형 구조상 공격적인 작업이 수반되고,\n"
   "이후 순위 하락 리스크를 다시 떠안게 되는 구조입니다.",
   Pt(12), False, RGBColor(0xCC,0xCC,0xDD))

# 내부 로직 테이블
data13 = [
    ["내부 로직 항목", "세부 내용"],
    ["키워드 세분화", "10유입 단위로 분할하여 패턴 감지를 회피"],
    ["유입경로 다양화", "매일 다른 키워드를 세팅하여 유입 경로를 다변화"],
    ["매체·미션 다양화", "단일 매체/미션 반복을 지양하고 다양한 조합으로 운영"],
]
tbl_shape13 = styled_table(s, ML, Cm(8.0), CW, 4, 2, data13, [Cm(8), Cm(20.867)])
# 다크 배경 테이블 텍스트 조정
tbl13 = tbl_shape13.table
for r in range(1, 4):
    for c in range(2):
        cell = tbl13.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x24, 0x37, 0x5E) if r % 2 == 1 else RGBColor(0x1E, 0x30, 0x55)
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = WHITE

tb(s, ML, Cm(12.2), CW, Cm(0.8),
   "※ 위 항목은 내부 로직의 일부 예시이며, 실제 적용 범위는 이보다 넓습니다.",
   Pt(10), False, RGBColor(0x99,0x99,0xAA))
add_footer_dark(s, 13)

# ====== SLIDE 14: 핵심 가치 ======
s = prs.slides.add_slide(blank)
add_top_line(s)
tb(s, ML, Cm(5.5), CW, Cm(1.5), "다른 업체들의 순위가 일제히 떨어질 때", Pt(22), False, DARK_TEXT, PP_ALIGN.CENTER)
tb(s, ML, Cm(7.5), CW, Cm(1.5), "귀사의 순위를 지킬 수 있는", Pt(28), True, DARK_TEXT, PP_ALIGN.CENTER)
tb(s, ML, Cm(9.5), CW, Cm(1.8), "가장 확실한 방법", Pt(38), True, BLUE, PP_ALIGN.CENTER)

# 골드 구분선
bar14 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(14.5), Cm(12.0), Cm(5), Cm(0.08))
bar14.fill.solid()
bar14.fill.fore_color.rgb = GOLD
bar14.line.fill.background()

tb(s, ML, Cm(13.0), CW, Cm(1.0),
   "안정적으로 유지되는 순위가 실질적인 매출을 만듭니다.", Pt(14), False, DARK_TEXT, PP_ALIGN.CENTER)
add_footer(s, 14)

# ====== SLIDE 15: AI로직연구소 ======
s = prs.slides.add_slide(blank)
set_bg(s, LIGHT_BG)
add_top_line(s)
add_logo(s)
add_title(s, "셀프마케팅 + AI로직연구소")
tb(s, ML, Cm(3.0), CW, Cm(0.7), "출시 예정 | 귀사도 바로 사용 가능", Pt(13), False, GRAY_TEXT)

data15 = [
    ["기능", "설명"],
    ["실시간 로직 모니터링", "경쟁사/동종업계 순위 및 추세를 3.5일 간격으로 감지하여 결과 제공"],
    ["맞춤 매체 추천", "현재 가장 안전하고 효율 좋은 트래픽 매체를 안내"],
    ["자동 등록 및 전환", "충전 시 자동으로 매체 등록 및 전환"],
]
styled_table(s, ML, Cm(4.2), Cm(16), 4, 2, data15, [Cm(5.5), Cm(10.5)])

placeholder_img(s, Cm(20), Cm(4.2), Cm(11.367), Cm(8.0), "이미지 자리: 솔루션 화면")

tb(s, ML, Cm(13.5), CW, Cm(2.0),
   "6장에서 보여드린 화면이 자체 운용 중인 솔루션입니다.\n"
   "미팅 자리에서 상세히 시연 및 설명드리겠습니다.",
   Pt(12), False, DARK_TEXT)
add_footer(s, 15)

# ====== SLIDE 16: 견적 ======
s = prs.slides.add_slide(blank)
add_top_line(s)
add_logo(s)
add_title(s, "견적 안내")
add_accent(s)

data16 = [
    ["지점", "현재순위", "보장형(첫달)", "보장형(2개월~)", "안전보장형(권장)"],
    ["마크원애비뉴점", "13위", "1,800,000원", "1,200,000원", "900,000원"],
    ["나성1호점", "15위", "1,900,000원", "1,200,000원", "900,000원"],
    ["나성2호점", "20위", "2,100,000원", "1,200,000원", "900,000원"],
    ["반곡점", "21위", "2,200,000원", "1,200,000원", "900,000원"],
    ["4지점 합계", "—", "8,000,000원", "4,800,000원", "3,600,000원"],
]
tbl16_shape = styled_table(s, ML, BODY_TOP, CW, 7, 5, data16,
    [Cm(6), Cm(4), Cm(6.289), Cm(6.289), Cm(6.289)])

# 안전보장형 열 골드 강조 + 합계행 처리
tbl16 = tbl16_shape.table
for r in range(1, 7):
    # 안전보장형 열 (index 4)
    cell = tbl16.cell(r, 4)
    cell.fill.solid()
    if r == 6:  # 합계행
        cell.fill.fore_color.rgb = GOLD
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.color.rgb = NAVY
            p.font.size = Pt(11)
    else:
        cell.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xEF)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True

    # 합계행 전체
    if r == 6:
        for c in range(5):
            cell_c = tbl16.cell(r, c)
            if c != 4:
                cell_c.fill.solid()
                cell_c.fill.fore_color.rgb = NAVY
                for p in cell_c.text_frame.paragraphs:
                    p.font.color.rgb = WHITE
                    p.font.bold = True
                    p.font.size = Pt(11)

tb(s, ML, Cm(11.0), CW, Cm(1.5),
   "※ 안전보장형은 4개 지점 동시 진행 시 지점당 900,000원 적용 (개별 진행 시 1,000,000원/지점)",
   Pt(11), False, GRAY_TEXT)
add_footer(s, 16)

# ====== SLIDE 17: 운영 조건 비교 ======
s = prs.slides.add_slide(blank)
set_bg(s, LIGHT_BG)
add_top_line(s)
add_logo(s)
add_title(s, "운영 조건 비교")
add_accent(s)

data17 = [
    ["항목", "보장형", "안전보장형(권장)"],
    ["작업 방식", "공격적 트래픽 투입", "내부 로직 중심 안정적 관리"],
    ["순위 상승 속도", "빠름", "점진적 / 안정적"],
    ["하락 리스크", "높음 (로직 변경 시 급락)", "낮음 (패턴 감지 회피)"],
    ["결제 방식", "노출 후 결제", "월 관리비"],
    ["최소 계약 기간", "없음", "없음"],
    ["리포팅 주기", "1주일", "1주일 (대시보드 상시 확인)"],
    ["월 비용 (4지점)", "첫달 800만 / 유지 480만", "360만원 (균일)"],
]
tbl17_shape = styled_table(s, ML, BODY_TOP, CW, 8, 3, data17, [Cm(6), Cm(11.434), Cm(11.433)])

# 하락 리스크 컬러코딩
tbl17 = tbl17_shape.table
# 보장형 하락리스크 = RED
cell_risk1 = tbl17.cell(3, 1)
for p in cell_risk1.text_frame.paragraphs:
    p.font.color.rgb = RED
    p.font.bold = True
# 안전보장형 하락리스크 = GREEN
cell_risk2 = tbl17.cell(3, 2)
for p in cell_risk2.text_frame.paragraphs:
    p.font.color.rgb = GREEN
    p.font.bold = True

# 안전보장형 열 배경
for r in range(1, 8):
    cell = tbl17.cell(r, 2)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xEF) if r % 2 == 1 else RGBColor(0xFD, 0xF4, 0xE2)

highlight_box(s, ML, Cm(12.5), CW, Cm(2.0),
    "온다마케팅은 안전보장형을 권장드립니다.\n"
    "최종 결정은 귀사의 상황과 판단에 따르며, 어떤 방식을 선택하시더라도 최선의 결과를 만들겠습니다.",
    RGBColor(0xEF,0xF6,0xFF), BLUE, DARK_TEXT, Pt(12))
add_footer(s, 17)

# ====== SLIDE 18: 마무리 ======
s = prs.slides.add_slide(blank)
set_bg(s, NAVY)
add_top_line(s, GOLD)

if os.path.exists(LOGO_PATH):
    s.shapes.add_picture(LOGO_PATH, Cm(12), Cm(3.0), Cm(10), Cm(3))

tb(s, ML, Cm(7.5), CW, Cm(2.0), "감사합니다", Pt(42), True, WHITE, PP_ALIGN.CENTER)

bar18 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(14.5), Cm(10.5), Cm(5), Cm(0.08))
bar18.fill.solid()
bar18.fill.fore_color.rgb = GOLD
bar18.line.fill.background()

tb(s, ML, Cm(11.5), CW, Cm(1.2),
   "어떤 방식을 선택하시더라도 최선의 결과를 만들겠습니다", Pt(16), False, GOLD, PP_ALIGN.CENTER)
tb(s, ML, Cm(14.0), CW, Cm(0.8), "온다마케팅", Pt(13), False, RGBColor(0x99,0x99,0xAA), PP_ALIGN.CENTER)
add_footer_dark(s, 18)

# === 저장 ===
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), '마드모아젤헤어_제안서.pptx')
prs.save(output_path)
print(f"✅ PPT 생성 완료: {output_path}")
print(f"✅ 슬라이드 수: {len(prs.slides)}")
print(f"✅ 파일 크기: {os.path.getsize(output_path)/1024:.0f}KB")